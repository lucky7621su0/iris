# create_presentation.py
# Generates a 10-slide PowerPoint for:
# "Portfolio Optimization and Risk Analysis using Machine Learning and Streamlit"
# Presenter: Prajwal, Dept: AI and DS, Supervisor: Kumar
#
# Requirements:
#   pip install python-pptx matplotlib numpy pandas seaborn
#
# Usage:
#   python create_presentation.py
#
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import seaborn as sns
import os

OUT_PPTX = "Portfolio_Optimization_and_Risk_Analysis_BE_Project.pptx"
IMG_DIR = "ppt_images"
os.makedirs(IMG_DIR, exist_ok=True)

prs = Presentation()

# Helper to add a slide with title and bullet points
def add_bullet_slide(title_text, bullets, notes_text=None, image_path=None):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    title.text = title_text

    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = b
        else:
            p = tf.add_paragraph()
            p.text = b
        p.level = 0

    if notes_text:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes_text

    if image_path and os.path.exists(image_path):
        # add image on right side
        left = Inches(5.2)
        top = Inches(1.5)
        width = Inches(4.0)
        slide.shapes.add_picture(image_path, left, top, width=width)

    return slide

# 1. Title Slide
title_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_layout)
slide.shapes.title.text = "Portfolio Optimization and Risk Analysis\nusing Machine Learning and Streamlit"
subtitle = slide.placeholders[1]
subtitle.text = "Prajwal — AI and DS\nSupervisor: Kumar"

notes = "Greet the audience, introduce yourself and summarize the project in one sentence: an interactive Streamlit app that combines ML forecasting with risk-aware portfolio optimization."
slide.notes_slide.notes_text_frame.text = notes

# Utility: create sample charts/images
def make_time_series_image(path):
    np.random.seed(42)
    t = np.arange(0, 200)
    series = np.cumsum(np.random.normal(0.001, 0.02, size=t.size)) + 1.0
    plt.figure(figsize=(6,3))
    plt.plot(t, series, color="#1f77b4")
    plt.title("Sample Asset Price (Synthetic)")
    plt.xlabel("Time")
    plt.ylabel("Price")
    plt.tight_layout()
    plt.savefig(path, dpi=150)
    plt.close()

def make_correlation_heatmap(path):
    np.random.seed(1)
    data = np.random.normal(size=(500,6))
    corr = pd.DataFrame(data, columns=list("ABCDEF")).corr()
    plt.figure(figsize=(4,3))
    sns.heatmap(corr, annot=True, cmap="coolwarm", vmin=-1, vmax=1)
    plt.title("Correlation Heatmap (Synthetic)")
    plt.tight_layout()
    plt.savefig(path, dpi=150)
    plt.close()

def make_confusion_matrix_image(path):
    cm = np.array([[85,15],[12,88]])
    plt.figure(figsize=(3,3))
    sns.heatmap(cm, annot=True, fmt="d", cmap="Blues")
    plt.title("Confusion Matrix (Placeholder)")
    plt.xlabel("Predicted")
    plt.ylabel("Actual")
    plt.tight_layout()
    plt.savefig(path, dpi=150)
    plt.close()

def make_efficient_frontier(path):
    np.random.seed(0)
    risks = np.linspace(0.05, 0.25, 50)
    returns = 0.02 + 0.3 * risks - 0.5 * risks**2 + np.random.normal(0, 0.002, risks.shape)
    plt.figure(figsize=(5,3))
    plt.plot(risks, returns, marker='o', linestyle='-')
    plt.scatter([0.12], [0.08], color='red', label='Optimized Portf')
    plt.xlabel("Risk (Std Dev)")
    plt.ylabel("Expected Return")
    plt.title("Efficient Frontier (Synthetic)")
    plt.legend()
    plt.tight_layout()
    plt.savefig(path, dpi=150)
    plt.close()

def make_cumulative_returns(path):
    np.random.seed(7)
    dates = pd.date_range(start="2019-01-01", periods=500)
    baseline = np.cumprod(1 + np.random.normal(0.0004, 0.01, size=len(dates)))
    ml = np.cumprod(1 + np.random.normal(0.0006, 0.009, size=len(dates)))
    plt.figure(figsize=(6,3))
    plt.plot(dates, baseline, label="Baseline")
    plt.plot(dates, ml, label="ML-based")
    plt.title("Cumulative Returns (Backtest, Synthetic)")
    plt.legend()
    plt.tight_layout()
    plt.savefig(path, dpi=150)
    plt.close()

# create images
img_ts = os.path.join(IMG_DIR, "timeseries.png")
img_corr = os.path.join(IMG_DIR, "heatmap.png")
img_cm = os.path.join(IMG_DIR, "confusion.png")
img_efficient = os.path.join(IMG_DIR, "efficient.png")
img_cum = os.path.join(IMG_DIR, "cumulative.png")
make_time_series_image(img_ts)
make_correlation_heatmap(img_corr)
make_confusion_matrix_image(img_cm)
make_efficient_frontier(img_efficient)
make_cumulative_returns(img_cum)

# 2. Agenda
bullets = [
    "Motivation",
    "Objectives & contributions",
    "Data & feature engineering",
    "ML models & risk measures",
    "Optimization methods",
    "System architecture & demo",
    "Results & evaluation",
    "Conclusion & future work"
]
notes = "Quickly walk through the agenda and set expectations for the demo near the end."
add_bullet_slide("Agenda", bullets, notes_text=notes)

# 3. Motivation & Problem Statement
bullets = [
    "Investors face market volatility and must balance return vs risk",
    "Traditional mean-variance uses noisy estimates → unstable allocations",
    "Machine Learning can improve return/volatility estimation",
    "Streamlit enables interactive decision-making and demo"
]
notes = "Explain why better estimates matter — small changes in expected returns cause big allocation changes. State the need for an interactive tool."
add_bullet_slide("Motivation & Problem Statement", bullets, notes_text=notes, image_path=img_ts)

# 4. Objectives & Contributions
bullets = [
    "Build ML models to predict asset returns and volatility",
    "Integrate risk measures (Variance, VaR, CVaR) into optimization",
    "Implement interactive Streamlit app for visualization and backtesting",
    "Contributions: comparative ML analysis & multi-objective optimization"
]
notes = "Summarize objectives and highlight novel aspects of the work."
add_bullet_slide("Objectives & Contributions", bullets, notes_text=notes)

# 5. Data & Feature Engineering
bullets = [
    "Data sources: historical daily prices (synthetic example)",
    "Features: log-returns, rolling mean/volatility, momentum, RSI",
    "Preprocessing: missing-value handling, scaling, walk-forward CV"
]
notes = "Describe dataset size and why technical indicators and macro variables were included."
add_bullet_slide("Data & Feature Engineering", bullets, notes_text=notes, image_path=img_corr)

# 6. Machine Learning Models & Risk Measures
bullets = [
    "Models tested: OLS, Ridge, Random Forest, XGBoost, small MLP",
    "Targets: 1-day / 1-week returns and volatility forecasts",
    "Evaluation: MSE/MAE for returns, coverage for VaR/CVaR; explainability via SHAP"
]
notes = "Discuss training approach and evaluation metrics. Explain CVaR preference for tail risk."
add_bullet_slide("ML Models & Risk Measures", bullets, notes_text=notes, image_path=img_cm)

# 7. Portfolio Optimization Methods
bullets = [
    "Mean-Variance (Markowitz) using ML expected returns",
    "CVaR minimization for tail-risk-aware allocations (LP-based)",
    "Constraints: weight bounds, turnover penalty, no-short rule (optionally)",
    "Backtest: rolling rebalancing with transaction-cost model"
]
notes = "High-level view of optimization objectives and constraints. Mention trade-offs between methods."
add_bullet_slide("Portfolio Optimization Methods", bullets, notes_text=notes, image_path=img_efficient)

# 8. System Architecture & Streamlit Demo
bullets = [
    "Pipeline: Data ingestion → Feature engineering → ML predictions → Optimizer → Streamlit UI",
    "Key libs: pandas, scikit-learn, xgboost, cvxpy, plotly, streamlit",
    "Demo features: select assets, choose model & risk metric, run backtest, visualize allocations"
]
notes = "Walk through the architecture and point out demo highlights (interactive controls, plots)."
# create a simple architecture placeholder image (reuse efficient frontier as placeholder)
add_bullet_slide("System Architecture & Streamlit Demo", bullets, notes_text=notes, image_path=img_efficient)

# 9. Results & Evaluation
bullets = [
    "Key findings (synthetic placeholders):",
    " - ML-based forecasts improved Sharpe by ~12% vs historical mean",
    " - CVaR optimization reduced max drawdown by ~18% (trade-off: slightly lower return)",
    " - Regularization reduced portfolio turnover"
]
notes = "Walk through the most important metrics and show cumulative returns / drawdown plots."
add_bullet_slide("Results & Evaluation", bullets, notes_text=notes, image_path=img_cum)

# 10. Conclusion, Limitations & Future Work
bullets = [
    "Conclusion: ML + risk-aware optimization can improve risk-adjusted performance",
    "Limitations: overfitting risk, data-snooping, transaction costs, limited universe",
    "Future work: expand asset classes, transaction-cost-aware rebalancing, deploy Dockerized Streamlit app"
]
notes = "Summarize takeaways and propose next steps. Invite questions and provide contact info."
add_bullet_slide("Conclusion, Limitations & Future Work", bullets, notes_text=notes)

# Save PPTX
prs.save(OUT_PPTX)
print(f"Saved PowerPoint: {OUT_PPTX}")
print(f"Generated images saved in folder: {IMG_DIR}")
